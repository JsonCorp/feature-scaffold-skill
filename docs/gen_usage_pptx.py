"""feature-scaffold-skill 사용 가이드 PPTX 생성기.

실행: python docs/gen_usage_pptx.py
산출물: docs/feature-scaffold-usage.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- 테마 ----
NAVY = RGBColor(0x0F, 0x1B, 0x2D)
PANEL = RGBColor(0x1E, 0x29, 0x3B)
BLUE = RGBColor(0x4F, 0x8F, 0xF7)
SLATE = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE = RGBColor(0xE2, 0xE8, 0xF0)
GREEN = RGBColor(0x6E, 0xE7, 0xB7)
FONT = "Malgun Gothic"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)  # rounded handled separately
    shp.fill.solid()
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space=6, mono=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        if isinstance(line, str):
            line = [(line, {})]
        for seg, sty in line:
            r = p.add_run()
            r.text = seg
            r.font.name = MONO if (mono or sty.get("mono")) else FONT
            r.font.size = Pt(sty.get("size", 18))
            r.font.bold = sty.get("bold", False)
            r.font.color.rgb = sty.get("color", WHITE)
    return tb


def accent_bar(slide):
    box(slide, 0, 0, Inches(0.18), SH, fill=BLUE)


def header(slide, kicker, title):
    accent_bar(slide)
    text(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.4),
         [[(kicker, {"size": 14, "bold": True, "color": BLUE})]])
    text(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.8),
         [[(title, {"size": 30, "bold": True, "color": WHITE})]])


def code_panel(slide, l, t, w, h, lines, size=14):
    panel = box(slide, l, t, w, h, fill=PANEL)
    runs = []
    for ln in lines:
        if isinstance(ln, tuple):
            txt, col = ln
        else:
            txt, col = ln, CODE
        runs.append([(txt if txt else " ", {"mono": True, "size": size, "color": col})])
    text(slide, l + Inches(0.25), t + Inches(0.18), w - Inches(0.5), h - Inches(0.36),
         runs, space=2)


# ---- 1. 표지 ----
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, Inches(2.6), SW, Inches(0.06), fill=BLUE)
text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.0),
     [[("feature-scaffold-skill", {"size": 46, "bold": True, "color": WHITE})]])
text(s, Inches(0.6), Inches(2.75), Inches(12), Inches(0.6),
     [[("Android 클린 아키텍처 기능 골격 생성 스킬", {"size": 22, "color": SLATE})]])
text(s, Inches(0.6), Inches(3.35), Inches(12), Inches(0.6),
     [[("사용 가이드", {"size": 22, "bold": True, "color": BLUE})]])
text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
     [[("github.com/JsonCorp/feature-scaffold-skill", {"size": 14, "mono": True, "color": SLATE})]])

# ---- 2. 이게 뭔가요 ----
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "WHAT", "이게 뭔가요?")
text(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(3.5), [
    [("기능 이름 한 줄로", {"size": 20, "color": WHITE}),
     (" 클린 아키텍처 기능 모듈 전체를 생성합니다.", {"size": 20, "color": WHITE})],
    [("", {})],
    [("• Scaffold(골격)", {"size": 18, "bold": True, "color": BLUE}),
     (" : 채워 넣을 표준 뼈대를 자동으로 세웁니다.", {"size": 18, "color": SLATE})],
    [("• Boilerplate(반복코드)", {"size": 18, "bold": True, "color": BLUE}),
     (" : 매번 똑같은 11~13개 파일을 손으로 안 씁니다.", {"size": 18, "color": SLATE})],
    [("", {})],
    [("→ domain · data · ui 3레이어 + Hilt DI + 테스트까지 한 번에.",
      {"size": 18, "color": GREEN})],
], space=10)

# ---- 3. 무엇을 만들어 주나 ----
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "OUTPUT", "무엇을 만들어 주나")
code_panel(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(4.9), [
    ("feature/<name>/", BLUE),
    "├── domain/",
    "│   ├── model/<Name>.kt            도메인 모델 (data class)",
    "│   ├── repository/<Name>Repository.kt   인터페이스",
    "│   └── usecase/Get<Name>UseCase.kt      단일 책임 UseCase",
    "├── data/",
    "│   ├── dto / mapper / api          DTO·매핑·Retrofit",
    "│   ├── repository/<Name>RepositoryImpl.kt",
    "│   └── di/<Name>Module.kt          Hilt (@Provides·@Binds)",
    "├── ui/",
    "│   ├── <Name>UiState.kt            sealed (Loading/Success/Error)",
    "│   ├── <Name>ViewModel.kt          HiltViewModel + StateFlow",
    "│   └── <Name>Screen.kt             Composable",
    ("└── test/                          UseCase·ViewModel 테스트", GREEN),
], size=14)

# ---- 4. 설치 ----
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "INSTALL", "설치 — 3단계")
text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.4),
     [[("1) 레포 클론", {"size": 16, "bold": True, "color": BLUE})]])
code_panel(s, Inches(0.6), Inches(2.25), Inches(12.1), Inches(0.95), [
    "git clone https://github.com/JsonCorp/feature-scaffold-skill.git",
    "cd feature-scaffold-skill",
], size=14)
text(s, Inches(0.6), Inches(3.4), Inches(12), Inches(0.4),
     [[("2) 설치 스크립트 실행", {"size": 16, "bold": True, "color": BLUE})]])
code_panel(s, Inches(0.6), Inches(3.8), Inches(5.95), Inches(1.7), [
    ("# macOS / Linux / Git Bash", SLATE),
    "./install.sh            # 이 프로젝트",
    "./install.sh /path/proj # 지정",
    "./install.sh --global   # 전역",
], size=13)
code_panel(s, Inches(6.75), Inches(3.8), Inches(5.95), Inches(1.7), [
    ("# Windows PowerShell", SLATE),
    ".\\install.ps1",
    ".\\install.ps1 -Target C:\\proj",
    ".\\install.ps1 -Global",
], size=13)
text(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.6),
     [[("3) Claude Code에서 ", {"size": 16, "color": WHITE}),
       ("/feature-scaffold", {"size": 16, "mono": True, "bold": True, "color": GREEN}),
       (" 호출", {"size": 16, "color": WHITE})]])

# ---- 5. 사용법 ----
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "USAGE", "사용법")
code_panel(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(0.8), [
    ("/feature-scaffold <FeatureName>[:field1[:Type],field2[:Type],...]", WHITE),
], size=16)
text(s, Inches(0.6), Inches(3.0), Inches(12), Inches(0.4),
     [[("예시", {"size": 16, "bold": True, "color": BLUE})]])
code_panel(s, Inches(0.6), Inches(3.4), Inches(12.1), Inches(2.2), [
    ("/feature-scaffold Profile", GREEN),
    "   → 기본 필드 id, name 으로 생성",
    "",
    ("/feature-scaffold Profile:id,email,age:Int", GREEN),
    "   → 지정 필드 + 타입(Int) 반영",
    "",
    ("/feature-scaffold Article:id,title,publishedAt:Long,likeCount:Int", GREEN),
    "   → String/Long/Int 혼합",
], size=14)

# ---- 6. 필드 커스터마이징 ----
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "FIELDS", "필드 커스터마이징")
text(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(2.2), [
    [("• 콜론(:) 뒤에 ", {"size": 18, "color": WHITE}),
     ("이름 또는 이름:타입", {"size": 18, "mono": True, "color": GREEN}),
     ("을 쉼표로 나열", {"size": 18, "color": WHITE})],
    [("• 타입 생략 시 ", {"size": 18, "color": WHITE}),
     ("String", {"size": 18, "mono": True, "color": GREEN}),
     (" 으로 간주", {"size": 18, "color": WHITE})],
    [("• 필드를 모두 생략하면 기본값 ", {"size": 18, "color": WHITE}),
     ("id, name", {"size": 18, "mono": True, "color": GREEN})],
    [("• 같은 필드가 도메인 모델 · DTO · Mapper 에 일관 반영", {"size": 18, "color": WHITE})],
], space=12)
code_panel(s, Inches(0.6), Inches(4.6), Inches(12.1), Inches(1.9), [
    ("Article:id,title,author,publishedAt:Long,likeCount:Int,isBookmarked:Boolean", BLUE),
    "",
    "data class Article(",
    "    val id: String, val title: String, val author: String,",
    "    val publishedAt: Long, val likeCount: Int, val isBookmarked: Boolean,",
    ")",
], size=13)

# ---- 7. 베이스 패키지 ----
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "CONFIG", "베이스 패키지 — 한 곳만 바꾸면 끝")
text(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(1.0), [
    [("생성 코드의 패키지 루트는 ", {"size": 18, "color": WHITE}),
     ("config.md", {"size": 18, "mono": True, "color": GREEN}),
     (" 한 줄로 결정됩니다.", {"size": 18, "color": WHITE})],
])
code_panel(s, Inches(0.6), Inches(2.9), Inches(12.1), Inches(0.8), [
    (".claude/skills/feature-scaffold/config.md", SLATE),
    ("base_package: com.example.app", GREEN),
], size=15)
text(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(1.6), [
    [("다른 프로젝트로 포팅?", {"size": 18, "bold": True, "color": BLUE})],
    [("• 이 한 줄만 com.your.app 으로 변경", {"size": 18, "color": WHITE})],
    [("• 또는 호출 시 ", {"size": 18, "color": WHITE}),
     ("--pkg=com.your.app", {"size": 18, "mono": True, "color": GREEN}),
     (" 로 1회성 오버라이드", {"size": 18, "color": WHITE})],
], space=10)

# ---- 8. 핵심 컨벤션 ----
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "RULES", "핵심 컨벤션")
text(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(4.5), [
    [("• 의존성 방향", {"size": 18, "bold": True, "color": BLUE}),
     ("  ViewModel → UseCase 만. Repository 직접 주입 금지.", {"size": 18, "color": WHITE})],
    [("• 인터페이스/구현 분리", {"size": 18, "bold": True, "color": BLUE}),
     ("  인터페이스는 domain, 구현은 data.", {"size": 18, "color": WHITE})],
    [("• DTO 격리", {"size": 18, "bold": True, "color": BLUE}),
     ("  DTO는 ui까지 노출 안 함. Mapper에서 변환.", {"size": 18, "color": WHITE})],
    [("• 상태 표현", {"size": 18, "bold": True, "color": BLUE}),
     ("  StateFlow + sealed interface (Loading/Success/Error).", {"size": 18, "color": WHITE})],
    [("• 취소 안전성", {"size": 18, "bold": True, "color": BLUE}),
     ("  try/catch + CancellationException rethrow. runCatching 금지.", {"size": 18, "color": WHITE})],
    [("• DI 바인딩", {"size": 18, "bold": True, "color": BLUE}),
     ("  Api는 @Provides, Repository는 @Binds. 공통 Retrofit 주입.", {"size": 18, "color": WHITE})],
], space=12)

# ---- 9. 마무리 ----
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, Inches(2.5), SW, Inches(0.06), fill=BLUE)
text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.9),
     [[("두 줄이면 시작", {"size": 36, "bold": True, "color": WHITE})]])
code_panel(s, Inches(0.6), Inches(2.9), Inches(12.1), Inches(1.1), [
    "git clone https://github.com/JsonCorp/feature-scaffold-skill.git",
    "cd feature-scaffold-skill && ./install.sh --global",
], size=15)
text(s, Inches(0.6), Inches(4.4), Inches(12), Inches(2.0), [
    [("그다음 Claude Code에서:", {"size": 18, "color": SLATE})],
    [("/feature-scaffold Article:id,title,author,likeCount:Int",
      {"size": 20, "mono": True, "bold": True, "color": GREEN})],
    [("", {})],
    [("문서 · 릴리스 · 예시 → 레포 README 참고", {"size": 16, "color": SLATE})],
], space=10)

prs.save("docs/feature-scaffold-usage.pptx")
print("saved docs/feature-scaffold-usage.pptx |", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
